#!/usr/bin/env python3
"""Generate Squad Profiles PowerPoint from HTML content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Colors ──
BLUE = RGBColor(0x0F, 0x3C, 0xC9)
DEEP_BLUE = RGBColor(0x0A, 0x28, 0x85)
GOLD = RGBColor(0xEF, 0xA7, 0x3D)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
BLACK = RGBColor(0x14, 0x14, 0x14)
GREY = RGBColor(0x59, 0x59, 0x59)
GREY_MUTED = RGBColor(0xA6, 0xA6, 0xA6)
BORDER = RGBColor(0xE0, 0xE0, 0xE0)
LIGHT_TINT = RGBColor(0xE7, 0xEB, 0xF8)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DOT_COLORS = {
    'client': GOLD,
    'acn': PURPLE,
    'either': 'gradient',  # handled specially
}

FONT = 'Nunito'

# ── Layout constants (inches) ──
SL_W = 13.333
SL_H = 7.5
PAD_L = 0.52
PAD_R = 0.52
PAD_T = 0.37

CONTENT_W = SL_W - PAD_L - PAD_R

def add_dot(slide, left, top, size, dot_type):
    """Add a colored dot (oval). For 'either', uses gradient gold→purple."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.line.fill.background()
    if dot_type == 'either':
        fill = shape.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = GOLD
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = PURPLE
        fill.gradient_stops[1].position = 1.0
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = DOT_COLORS.get(dot_type, GREY)


def set_text(tf, text, size=8, bold=False, color=BLACK, italic=False, align=PP_ALIGN.LEFT):
    """Set single paragraph text in a text frame."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    run.font.italic = italic
    return p


def add_para(tf, text, size=8, bold=False, color=BLACK, italic=False, space_before=0):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    run.font.italic = italic
    if space_before:
        p.space_before = Pt(space_before)
    return p


def add_textbox(slide, left, top, width, height, text, size=8, bold=False, color=BLACK, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box with single text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Set margins tight
    txBox.text_frame.margin_left = Emu(0)
    txBox.text_frame.margin_right = Emu(0)
    txBox.text_frame.margin_top = Emu(0)
    txBox.text_frame.margin_bottom = Emu(0)
    set_text(tf, text, size, bold, color, italic, align)
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=0.75, border_dash=None, corner_radius=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
        if border_dash:
            shape.line.dash_style = border_dash
    else:
        shape.line.fill.background()
    # Adjust corner radius
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    else:
        shape.adjustments[0] = 0.04
    return shape


def add_role(slide, left, top, width, dot_type, name, desc, name_suffix=None):
    """Add a role row: dot + name + description."""
    dot_size = 0.065
    dot_left = left
    dot_top = top + 0.02
    add_dot(slide, dot_left, dot_top, dot_size, dot_type)

    text_left = left + 0.10
    text_width = width - 0.10

    txBox = slide.shapes.add_textbox(Inches(text_left), Inches(top), Inches(text_width), Inches(0.50))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    # Name paragraph
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = BLACK
    run.font.name = FONT
    if name_suffix:
        run2 = p.add_run()
        run2.text = f"  {name_suffix}"
        run2.font.size = Pt(7)
        run2.font.bold = True
        run2.font.color.rgb = BLUE
        run2.font.name = FONT

    # Desc paragraph
    p2 = tf.add_paragraph()
    p2.space_before = Pt(1)
    run3 = p2.add_run()
    run3.text = desc
    run3.font.size = Pt(6.5)
    run3.font.color.rgb = GREY
    run3.font.name = FONT

    return txBox


# ── Build the presentation ──
prs = Presentation()
prs.slide_width = Inches(SL_W)
prs.slide_height = Inches(SL_H)

slide_layout = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(slide_layout)

# ── Header ──
add_textbox(slide, PAD_L, PAD_T, 4, 0.20, "TALENT & TEAM MODEL", size=9, bold=False, color=GREY)
add_textbox(slide, PAD_L, PAD_T + 0.22, 6, 0.35, "Squad Profiles", size=21, bold=True, color=BLACK)

# ── Legend (top right) ──
legend_x = SL_W - PAD_R - 2.0
legend_y = PAD_T
legend_items = [
    ('client', 'Typically client filled'),
    ('acn', 'Typically Accenture filled'),
    ('either', 'Filled by either'),
]
for i, (dot_type, label) in enumerate(legend_items):
    y = legend_y + i * 0.22
    add_dot(slide, legend_x, y + 0.02, 0.09, dot_type)
    add_textbox(slide, legend_x + 0.14, y, 1.8, 0.18, label, size=8, color=GREY)

# ── Portfolio bar ──
port_y = PAD_T + 0.70
port_h = 0.42
port_items = [
    (None, 'Portfolio', '', 'acn'),
    ('acn', 'Delivery Leadership', 'Velocity and burn across squads; manages dependencies and escalation', None),
    ('either', 'VRO & Adoption', 'Tracks agent adoption, cost-per-task, value realization vs business case', None),
    ('either', 'Enterprise & Strategic Architecture', 'Aligns agent platform with enterprise AI strategy and integration landscape', None),
    ('client', 'Security & Privacy', 'Responsible AI governance, data classification, model risk approvals', None),
]

# Blue "Portfolio" label
port_label_w = 0.85
add_rect(slide, PAD_L, port_y, port_label_w, port_h, fill_color=BLUE, corner_radius=0.06)
add_textbox(slide, PAD_L + 0.08, port_y + 0.10, port_label_w - 0.16, 0.22, "Portfolio", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Remaining portfolio items
port_item_w = (CONTENT_W - port_label_w - 0.04) / 4
for i, (dot_type, title, desc, _) in enumerate(port_items[1:], 0):
    x = PAD_L + port_label_w + 0.04 + i * (port_item_w + 0.01)
    # Background rect
    add_rect(slide, x, port_y, port_item_w, port_h, fill_color=WHITE, border_color=BORDER, border_width=0.5, corner_radius=0.03)
    # Dot
    add_dot(slide, x + 0.08, port_y + 0.13, 0.07, dot_type)
    # Title
    add_textbox(slide, x + 0.20, port_y + 0.06, port_item_w - 0.28, 0.16, title, size=8, bold=True, color=BLACK)
    # Desc
    if desc:
        add_textbox(slide, x + 0.20, port_y + 0.22, port_item_w - 0.28, 0.18, desc, size=6, color=GREY)

# ── Content grid ──
grid_y = port_y + port_h + 0.15
grid_h = SL_H - grid_y - 0.50  # Leave room for footer

# Persistent squad panel
pers_w = 2.10
pers_x = PAD_L
add_rect(slide, pers_x, grid_y, pers_w, grid_h, fill_color=LIGHT_TINT, corner_radius=0.04)

# Persistent header
add_textbox(slide, pers_x + 0.10, grid_y + 0.08, pers_w - 0.20, 0.18, "Persistent Squad", size=9, bold=True, color=BLUE)
add_textbox(slide, pers_x + 0.10, grid_y + 0.26, pers_w - 0.20, 0.14, "Present across all phases", size=6.5, color=GREY, italic=True)

# Persistent roles
persistent_roles = [
    ('either', 'Squad Lead', 'Manages sprint planning, capacity allocation and cross-role dependencies; escalation point within the squad', None),
    ('acn', 'AI Solution Architect', 'Sits cross-squad; defines agent architecture, ensures alignment on tool graph, memory and guardrails', None),
    ('client', 'Product Owner', 'Prioritises features by business value and technical feasibility; defines evaluation acceptance criteria', None),
    ('acn', 'Forward Deployed AI Engineers', 'Design and implement agent behaviors, tool integrations, evaluation pipelines and orchestration layers', '2-3x'),
    ('acn', 'Agent Behaviour Analyst', 'Translates business processes into agent behavior specifications and eval criteria; anticipates how agent capabilities replace static SOPs', None),
    ('either', 'Business SME', 'Supplies domain ground-truth and validates agent outputs against real-world accuracy and compliance standards', None),
]

role_y = grid_y + 0.46
for dot_type, name, desc, suffix in persistent_roles:
    add_role(slide, pers_x + 0.10, role_y, pers_w - 0.20, dot_type, name, desc, suffix)
    role_y += 0.58

# ── Phase columns ──
phases_x = pers_x + pers_w + 0.10
phases_w = CONTENT_W - pers_w - 0.10
phase_gap = 0.08
phase_w = (phases_w - phase_gap * 3) / 4

# Phase data
phases = [
    {
        'title': 'Discovery',
        'sub': 'Validate agent feasibility, define target journeys and establish evaluation criteria',
        'roles': [
            ('either', 'AI Solution Architect', 'Maps source data, APIs and integration points; validates what agents can feasibly access and act on', None),
            ('acn', 'UX Designer', 'Produces high-fidelity prototypes to validate user interaction patterns with agent-driven interfaces', None),
            ('acn', 'Human-AI Experience (HAX) Designer', 'Designs agent persona, trust patterns, escalation moments and multi-modal interaction across voice, chat and gesture', None),
        ],
        'note': 'FD AI Engineers at 2x during Discovery (bootstrapping agents and evals)',
        'dashed': False,
    },
    {
        'title': 'Build & Iterate',
        'sub': 'Incremental agent delivery with continuous evaluation and prompt refinement',
        'roles': [
            ('acn', 'Platform / Infra Engineer', 'Edge/core inference split, token metering, DFC integration', None),
            ('acn', 'Data Engineer', 'LifeGraph pipelines, cluster embeddings, signal infrastructure', None),
            ('acn', 'Front-end Developers', 'Builds native applications, agent-facing interfaces, streaming chat components and operational dashboards', '2x'),
            ('acn', 'QA & Eval Engineer', 'Builds evaluation harnesses for non-deterministic agent behavior and regression suites for deterministic components', None),
        ],
        'note': 'FD AI Engineers at peak 3x; observability, eval coverage and agent performance',
        'dashed': False,
    },
    {
        'title': 'Harden & Launch',
        'sub': 'Adversarial testing, guardrail validation and production certification',
        'roles': [
            ('client', 'Security Architect', 'Conducts security assessment including LLM-specific risks: prompt injection, data leakage, model exfiltration', None),
            ('acn', 'DevOps / AIOps Engineer', 'Implements CI/CD for agent deployments and eval suites; manages canary releases, rollback and production observability', None),
            ('acn', 'Site Reliability Engineer', 'Defines and enforces SLOs for inference latency, agent accuracy and availability; incident response for model degradation', None),
            ('acn', 'Technical & AI QA', 'Systematically tests agent edge-cases through adversarial scenarios; validates guardrail effectiveness at scale', None),
            ('client', 'Security Representative', 'Responsible AI and security signoff before each release tier', None),
            ('client', 'Jio Integration Lead', 'Manages API integration with JioMart, JioFinance, JioCinema and UPI gateway; validates end-to-end data flows', None),
        ],
        'note': 'FD AI Engineers at 2x; hardening agent behavior and production readiness',
        'dashed': False,
    },
    {
        'title': 'Steady-State',
        'sub': 'Smaller feature cycles; model governance, agent refinement and new signals',
        'roles': [
            ('either', 'Squad Lead', 'Manages sprint cadence across smaller feature cycles', None),
            ('acn', 'AI Solution Architect', 'Cross-squad alignment; architecture governance as new signals land', None),
            ('client', 'Product Owner', 'Prioritises new features and signals; evaluates agent performance', None),
            ('acn', 'FD AI Engineers', 'Prompt refinement, eval drift detection, model version transitions', '1-2x'),
            ('acn', 'Platform / Infra Engineer', 'Autoscaling, token cost governance, inference latency monitoring', None),
            ('acn', 'Data Engineer', 'New signal onboarding, embedding freshness, feedback loop pipelines', None),
        ],
        'note': None,
        'dashed': True,
        'preamble': 'AI-native products do not transition to traditional BAU. Agents learn, models drift, new signals emerge. The squad stays on at 25-40% capacity running shorter feature cycles.',
    },
]

for i, phase in enumerate(phases):
    px = phases_x + i * (phase_w + phase_gap)

    # Column background
    fill = OFF_WHITE if phase.get('dashed') else WHITE
    border_dash = 4 if phase.get('dashed') else None  # MSO_LINE_DASH.DASH = 4
    rect = add_rect(slide, px, grid_y, phase_w, grid_h, fill_color=fill, border_color=BORDER, border_width=0.75, corner_radius=0.04)
    if border_dash:
        rect.line.dash_style = border_dash

    # Phase header background line
    header_h = 0.52

    # Title
    add_textbox(slide, px + 0.10, grid_y + 0.08, phase_w - 0.20, 0.18, phase['title'], size=9, bold=True, color=BLACK)
    # Subtitle
    add_textbox(slide, px + 0.10, grid_y + 0.26, phase_w - 0.20, 0.26, phase['sub'], size=6, color=GREY, italic=True)

    # Roles
    ry = grid_y + header_h + 0.06

    # Steady-state preamble
    if phase.get('preamble'):
        tb = add_textbox(slide, px + 0.10, ry, phase_w - 0.20, 0.55, '', size=7, color=GREY)
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        # Add text with bold "25-40% capacity"
        parts = phase['preamble'].split('25-40% capacity')
        run1 = p.add_run()
        run1.text = parts[0]
        run1.font.size = Pt(7)
        run1.font.color.rgb = GREY
        run1.font.name = FONT
        run_bold = p.add_run()
        run_bold.text = '25-40% capacity'
        run_bold.font.size = Pt(7)
        run_bold.font.bold = True
        run_bold.font.color.rgb = BLACK
        run_bold.font.name = FONT
        run2 = p.add_run()
        run2.text = parts[1] if len(parts) > 1 else ''
        run2.font.size = Pt(7)
        run2.font.color.rgb = GREY
        run2.font.name = FONT
        ry += 0.55

    for dot_type, name, desc, suffix in phase['roles']:
        add_role(slide, px + 0.10, ry, phase_w - 0.20, dot_type, name, desc, suffix)
        # Smaller spacing for phases with many roles
        if len(phase['roles']) > 4:
            ry += 0.48
        else:
            ry += 0.55

    # Phase note
    if phase.get('note'):
        note_y = grid_y + grid_h - 0.28
        # Note background
        add_rect(slide, px + 0.01, note_y, phase_w - 0.02, 0.26, fill_color=OFF_WHITE, border_color=BORDER, border_width=0.5, corner_radius=0.02)
        add_textbox(slide, px + 0.10, note_y + 0.04, phase_w - 0.20, 0.20, phase['note'], size=6, color=GREY, italic=True)

# ── Footer ──
footer_y = SL_H - 0.35

# Left
txBox = slide.shapes.add_textbox(Inches(PAD_L), Inches(footer_y), Inches(4), Inches(0.20))
tf = txBox.text_frame
tf.margin_left = Emu(0)
tf.margin_top = Emu(0)
p = tf.paragraphs[0]
for text, bold in [("Reliance", True), (" x ", False), ("Accenture", True), (" Strategic Partnership", False)]:
    run = p.add_run()
    run.text = text
    run.font.size = Pt(7.5)
    run.font.bold = bold
    run.font.color.rgb = GREY_MUTED
    run.font.name = FONT

# Right
add_textbox(slide, SL_W - PAD_R - 3.5, footer_y, 3.5, 0.20, "Copyright Accenture 2026. All rights reserved.", size=7.5, color=GREY_MUTED, align=PP_ALIGN.RIGHT)

# ── Save ──
output = '/Users/v.eguren/Documents/GitHub/jio-demos/presentations/squad-profiles.pptx'
prs.save(output)
print(f"Saved: {output}")
